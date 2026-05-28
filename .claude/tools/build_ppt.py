#!/usr/bin/env python3
"""
Odoo 工時表同步與 PPT 自動生成腳本
自動動態計算「執行當天前兩週的週一至週五」，建立目錄並生成左右雙表格 PPT。
"""

import os
import sys
import requests
from datetime import datetime, timedelta
from bs4 import BeautifulSoup
from pathlib import Path
from collections import defaultdict

# 匯入 PPT 相關套件
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def remove_images_only(html_str):
    if not html_str or not isinstance(html_str, str):
        return ""
    soup = BeautifulSoup(html_str, "html.parser")
    for img in soup.find_all("img"):
        img.decompose()
    return str(soup)

def get_previous_two_weeks_range():
    """
    動態計算執行當天的「前兩週（上上週與上週）」的週一與週五日期。
    不論在週幾觸發，都能精確對齊。
    """
    today = datetime.now().date()
    
    # 1. 找出本週週一 (today.weekday() 週一是 0, 週日是 6)
    this_week_monday = today - timedelta(days=today.weekday())
    
    # 2. 觸發日的前兩週週一 (上上週一)
    two_weeks_ago_monday = this_week_monday - timedelta(weeks=2)
    
    # 3. 觸發日的前一週週五 (上週五) -> 上上週一往後推 11 天就是上週五
    last_week_friday = two_weeks_ago_monday + timedelta(days=11)
    
    return two_weeks_ago_monday.strftime("%Y-%m-%d"), last_week_friday.strftime("%Y-%m-%d")

def generate_ppt(start_date, end_date, project_stats, total_all_hours, start_path):
    """
    生成包含左右雙表格的 PPT 簡報
    """
    # 轉換日期格式為 MM/DD
    dt_start = datetime.strptime(start_date, "%Y-%m-%d").strftime("%m/%d")
    dt_end = datetime.strptime(end_date, "%Y-%m-%d").strftime("%m/%d")
    
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # 使用完全空白頁面
    slide = prs.slides.add_slide(blank_layout)
    
    # 1. 設定大標題 (動態帶入產出的前兩週時間)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"過去兩周工作進度 {dt_start}~{dt_end}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Microsoft JhengHei"
    
    # 準備左邊表格數據 (按工時高低排序)
    left_data = []
    for proj, data in sorted(project_stats.items(), key=lambda x: x[1]["hours"], reverse=True):
        proj_hours = data["hours"]
        task_count = len(data["tasks"])
        percentage = (proj_hours / total_all_hours * 100) if total_all_hours > 0 else 0.0
        
        hours_str = f"{proj_hours:.1f}hr ({percentage:.1f}%)"
        left_data.append([proj, hours_str, str(task_count)])
        
    # 2. 繪製左邊表格 (過去兩周工作進度)
    left_rows = len(left_data) + 1  # 資料列 + 1 標頭列
    left_cols = 3
    left_table_shape = slide.shapes.add_table(
        left_rows, left_cols, Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.35 * left_rows)
    )
    left_table = left_table_shape.table
    
    # 設定左表欄寬
    left_table.columns[0].width = Inches(2.2) # 項目名稱
    left_table.columns[1].width = Inches(1.5) # 總工時(hr)
    left_table.columns[2].width = Inches(0.8) # 單號數量
    
    # 左表標頭
    left_headers = ["項目名稱", "總工時(hr)", "單號數量"]
    for col_idx, text in enumerate(left_headers):
        cell = left_table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # 專業藍
        
    # 左表填入資料
    for row_idx, row_data in enumerate(left_data):
        for col_idx, text in enumerate(row_data):
            cell = left_table.cell(row_idx + 1, col_idx)
            cell.text = text

    # 3. 繪製右邊表格 (未來兩周工作計畫 - 留空供後續填寫)
    right_rows = max(6, left_rows)  # 至少保留 6 列空間
    right_cols = 2
    right_table_shape = slide.shapes.add_table(
        right_rows, right_cols, Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.35 * right_rows)
    )
    right_table = right_table_shape.table
    
    # 設定右表欄寬
    right_table.columns[0].width = Inches(2.7) # 未來兩周工作計畫項目名稱
    right_table.columns[1].width = Inches(1.5) # 預估工時
    
    # 右表標頭
    right_headers = ["未來兩周工作計畫項目名稱", "預估工時"]
    for col_idx, text in enumerate(right_headers):
        cell = right_table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(128, 128, 128)  # 質感灰
        
    # 統一美化表格文字與對齊
    for table in [left_table, right_table]:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "Microsoft JhengHei"
                    paragraph.font.size = Pt(11)
                    paragraph.alignment = PP_ALIGN.CENTER
                    
    # 儲存簡報
    ppt_file = start_path / "工作週報.pptx"
    prs.save(ppt_file)
    print(f"[SUCCESS] PPT 簡報已成功動態生成至: {ppt_file}")

def main():
    if len(sys.argv) < 8:
        print("[ERROR] 參數不足。用法: python build_ppt.py <URL> <DB> <USER> <PWD> <USER_ID> <START_DIR> <PREFIX> [SKIP_IDS]")
        sys.exit(1)

    ODOO_URL = sys.argv[1]
    DB_NAME = sys.argv[2]
    USERNAME = sys.argv[3]
    PASSWORD = sys.argv[4]
    USER_ID = int(sys.argv[5])
    START_DIR = sys.argv[6]
    PREFIX = sys.argv[7]
    SKIP_IDS = set(sys.argv[8].split(",")) if len(sys.argv) > 8 and sys.argv[8] else set()

    CHROME_UA = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"

    session = requests.Session()
    session.headers.update({"User-Agent": CHROME_UA})

    # Odoo 登入
    auth_url = f"{ODOO_URL}/web/session/authenticate"
    auth_payload = {
        "jsonrpc": "2.0",
        "method": "call",
        "params": {"db": DB_NAME, "login": USERNAME, "password": PASSWORD}
    }

    try:
        auth_resp = session.post(auth_url, json=auth_payload).json()
        if "error" in auth_resp:
            print(f"[ERROR] Odoo 登入失敗: {auth_resp['error']}")
            sys.exit(1)
    except Exception as e:
        print(f"[ERROR] 連線失敗: {e}")
        sys.exit(1)

    # 4. 【核心改動】動態獲取前兩週週一至週五的時間
    start_date, end_date = get_previous_two_weeks_range()
    print(f"[INFO] 系統自動判定觸發日前兩週區間: {start_date} ~ {end_date}")

    # 獲取工時表列表
    call_url = f"{ODOO_URL}/web/dataset/call_kw"
    timesheet_payload = {
        "jsonrpc": "2.0",
        "method": "call",
        "params": {
            "model": "account.analytic.line",
            "method": "search_read",
            "args": [],
            "kwargs": {
                "domain": [
                    ["user_id", "=", USER_ID],
                    ["date", ">=", start_date],
                    ["date", "<=", end_date]
                ],
                "fields": ["id", "name", "date", "unit_amount", "project_id", "task_id"],
                "order": "date desc",
                "limit": 150
            }
        }
    }

    ts_resp = session.post(call_url, json=timesheet_payload).json()
    timesheets = ts_resp.get("result", [])

    start_path = Path(START_DIR)
    start_path.mkdir(parents=True, exist_ok=True)

    if not timesheets:
        print(f"[INFO] 區間 {start_date} ~ {end_date} 內沒有任何工時紀錄，跳過檔案與 PPT 生成。")
        sys.exit(0)

    project_stats = defaultdict(lambda: {"hours": 0.0, "tasks": set()})
    total_all_hours = 0.0

    # 處理工時資料與建立目錄
    for ts in timesheets:
        ts_id = ts.get("id")
        ts_date = ts.get("date", "未知日期")
        
        # 嚴格過濾掉週六日
        if datetime.strptime(ts_date, "%Y-%m-%d").date().weekday() in (5, 6):
            continue

        ts_hours = ts.get("unit_amount", 0.0)
        ts_description = remove_images_only(ts.get("name")).strip()
        project_name = ts["project_id"][1] if ts.get("project_id") else "未知專案"
        task_name = ts["task_id"][1] if ts.get("task_id") else "未關聯任務"

        project_stats[project_name]["hours"] += ts_hours
        if ts_description:
            project_stats[project_name]["tasks"].add(ts_description)
        total_all_hours += ts_hours

        if str(ts_id) in SKIP_IDS:
            continue

        ts_dir = start_path / f"{PREFIX}{ts_id}"
        if ts_dir.exists():
            continue

        ts_dir.mkdir(parents=True, exist_ok=True)
        original_content = f"---id---\n{ts_id}\n---date---\n{ts_date}\n---hours---\n{ts_hours}\n---project---\n{project_name}\n---task---\n{task_name}\n---description---\n{ts_description}"
        (ts_dir / "original.txt").write_text(original_content, encoding="utf-8")

    # 執行動態生成 PPT 簡報
    generate_ppt(start_date, end_date, project_stats, total_all_hours, start_path)

if __name__ == "__main__":
    main()
